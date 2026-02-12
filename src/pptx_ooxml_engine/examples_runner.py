from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation

from .engine import generate_pptx


def _create_demo_template(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Demo Template Title"
    prs.save(str(path))


def _create_demo_reuse_library(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Demo Reuse Slide"
    prs.save(str(path))


def _load_example_ops(example_path: Path) -> dict:
    return json.loads(example_path.read_text(encoding="utf-8"))


def _copy_ops_available() -> bool:
    try:
        from .engine import _import_copy_ops

        _import_copy_ops()
        return True
    except ModuleNotFoundError:
        return False


def _uses_copy_op(ops_doc: dict) -> bool:
    operations = ops_doc.get("operations", [])
    return any(isinstance(item, dict) and item.get("op") == "copy_slide" for item in operations)


def generate_example_outputs(output_dir: str | Path) -> list[Path]:
    out_dir = Path(output_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    assets_dir = out_dir / "_assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    template_pptx = assets_dir / "template_demo.pptx"
    reuse_lib_pptx = assets_dir / "reuse_demo.pptx"
    _create_demo_template(template_pptx)
    _create_demo_reuse_library(reuse_lib_pptx)

    package_root = Path(__file__).resolve().parents[2]
    ops_dir = package_root / "examples" / "ops"
    ops_files = [
        ops_dir / "01_template_only_generate.json",
        ops_dir / "02_template_plus_reuse_library.json",
    ]

    copy_available = _copy_ops_available()
    generated: list[Path] = []
    for idx, ops_file in enumerate(ops_files, start=1):
        ops_doc = _load_example_ops(ops_file)
        if _uses_copy_op(ops_doc) and not copy_available:
            continue
        ops_doc["template_pptx"] = str(template_pptx)
        libs = list(ops_doc.get("reuse_slide_libraries", []))
        if libs:
            ops_doc["reuse_slide_libraries"] = [str(reuse_lib_pptx)]
        output_path = out_dir / f"example_{idx:02d}.pptx"
        result = generate_pptx(
            template_pptx=template_pptx,
            ops=ops_doc,
            output_pptx=output_path,
            verify=True,
        )
        generated.append(result.output_path)
    if not generated:
        raise RuntimeError("No examples generated. Install pptx-copy-ops or provide non-copy examples.")
    return generated


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Run bundled ops examples and generate demo PPTX outputs.")
    parser.add_argument("--output-dir", required=True, help="Output directory for generated example PPTX files.")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = _build_parser()
    args = parser.parse_args(argv)
    outputs = generate_example_outputs(args.output_dir)
    for output in outputs:
        print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
