from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.slide import Slide
from pptx.util import Inches, Pt

from .models import (
    AddImageOp,
    AddShapeOp,
    AddTableOp,
    AddTextBoxOp,
    AlignShapesOp,
    CopyMode,
    CopySlideOp,
    CreateSlideOnLayoutOp,
    DeleteSlideOp,
    DistributeShapesOp,
    MoveSlideOp,
    Operation,
    OperationPlan,
    ParagraphSpec,
    RewriteTextOp,
    SetNotesOp,
    SetShapeTextOp,
    SetSlideBackgroundOp,
    SetSlideLayoutOp,
    SetSlideSizeOp,
    parse_plan,
    parse_ops,
)
from .verify import verify_pptx

_SLIDE_LAYOUT_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
_BULLET_TAGS = (qn("a:buNone"), qn("a:buChar"), qn("a:buAutoNum"))
_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}
_VERTICAL_ANCHOR_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}
_SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
}


def _import_copy_ops():
    try:
        from pptx_copy_ops import SlideCopier, SlideSpec  # type: ignore

        return SlideCopier, SlideSpec
    except ModuleNotFoundError:
        pass

    workspace_root = Path(__file__).resolve().parents[3]
    fallback_src = workspace_root / "pptx-copy-ops" / "src"
    if fallback_src.exists() and str(fallback_src) not in sys.path:
        sys.path.insert(0, str(fallback_src))
    try:
        from pptx_copy_ops import SlideCopier, SlideSpec  # type: ignore

        return SlideCopier, SlideSpec
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "pptx-copy-ops is required for copy_slide operations. "
            "Install it or provide sibling path `pptx-copy-ops/src`."
        ) from exc


@dataclass
class ApplyResult:
    output_path: Path
    operations_applied: int
    verify_issues: list[str]


def _iter_text_shapes(slide: Slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape


def _hex_to_rgb(color_hex: str) -> RGBColor:
    return RGBColor.from_string(color_hex.lstrip("#"))


def _slide_or_raise(presentation: Presentation, slide_index: int, op_name: str):
    if slide_index >= len(presentation.slides):
        raise IndexError(f"{op_name} slide_index out of range: {slide_index}, total={len(presentation.slides)}")
    return presentation.slides[slide_index]


def _shape_by_name(slide: Slide, shape_name: str):
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    raise ValueError(f"shape not found by name: {shape_name}")


def _shape_by_index(slide: Slide, shape_index: int):
    if shape_index >= len(slide.shapes):
        raise IndexError(f"shape_index out of range: {shape_index}, total={len(slide.shapes)}")
    return slide.shapes[shape_index]


def _set_paragraph_list_style(paragraph, list_type: str) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        if child.tag in _BULLET_TAGS:
            p_pr.remove(child)
    if list_type == "bullet":
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        p_pr.append(bullet)
    elif list_type == "number":
        auto_num = OxmlElement("a:buAutoNum")
        auto_num.set("type", "arabicPeriod")
        auto_num.set("startAt", "1")
        p_pr.append(auto_num)
    else:
        bu_none = OxmlElement("a:buNone")
        p_pr.append(bu_none)


def _apply_paragraph_style(paragraph, spec: ParagraphSpec) -> None:
    paragraph.level = spec.level
    _set_paragraph_list_style(paragraph, spec.list_type)
    if spec.alignment is not None:
        paragraph.alignment = _ALIGN_MAP[spec.alignment]
    if spec.line_spacing is not None:
        paragraph.line_spacing = spec.line_spacing
    if spec.space_before_pt is not None:
        paragraph.space_before = Pt(spec.space_before_pt)
    if spec.space_after_pt is not None:
        paragraph.space_after = Pt(spec.space_after_pt)

    font = paragraph.font
    if spec.font_size_pt is not None:
        font.size = Pt(spec.font_size_pt)
    if spec.bold is not None:
        font.bold = spec.bold
    if spec.italic is not None:
        font.italic = spec.italic
    if spec.color_hex is not None:
        font.color.rgb = _hex_to_rgb(spec.color_hex)


def _write_text_frame(
    text_frame,
    text: str | None,
    paragraphs: list[ParagraphSpec],
    vertical_anchor: str | None,
    word_wrap: bool | None,
) -> None:
    if vertical_anchor is not None:
        text_frame.vertical_anchor = _VERTICAL_ANCHOR_MAP[vertical_anchor]
    if word_wrap is not None:
        text_frame.word_wrap = word_wrap

    payload = paragraphs if paragraphs else [ParagraphSpec(text=text or "")]
    text_frame.clear()
    for idx, spec in enumerate(payload):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = spec.text
        _apply_paragraph_style(paragraph, spec)


def _set_slide_body(slide: Slide, body: str) -> None:
    for shape in _iter_text_shapes(slide):
        if getattr(shape, "is_placeholder", False):
            try:
                placeholder_type = shape.placeholder_format.type
            except Exception:
                placeholder_type = None
            # BODY placeholder
            if placeholder_type == 2:
                shape.text_frame.text = body
                return
    for shape in _iter_text_shapes(slide):
        if shape != slide.shapes.title:
            shape.text_frame.text = body
            return


def _apply_rewrite(op: RewriteTextOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "rewrite_text")
    replaced = 0
    for shape in _iter_text_shapes(slide):
        if op.shape_name and shape.name != op.shape_name:
            continue
        text = shape.text_frame.text
        if op.find not in text:
            continue
        count = 1 if op.occurrence == "first" else -1
        shape.text_frame.text = text.replace(op.find, op.replace, count)
        replaced += 1
        if op.occurrence == "first":
            break
    if replaced == 0:
        raise ValueError(f"rewrite_text cannot find target text on slide {op.slide_index}: {op.find!r}")


def _apply_create(op: CreateSlideOnLayoutOp, presentation: Presentation) -> None:
    if op.layout_index >= len(presentation.slide_layouts):
        raise IndexError(f"layout_index out of range: {op.layout_index}, total={len(presentation.slide_layouts)}")
    slide = presentation.slides.add_slide(presentation.slide_layouts[op.layout_index])
    if op.title and slide.shapes.title is not None:
        slide.shapes.title.text = op.title
    if op.body:
        _set_slide_body(slide, op.body)


def _apply_delete(op: DeleteSlideOp, presentation: Presentation) -> None:
    if op.slide_index >= len(presentation.slides):
        raise IndexError(f"delete_slide slide_index out of range: {op.slide_index}, total={len(presentation.slides)}")
    r_id = presentation.slides._sldIdLst[op.slide_index].rId
    presentation.part.drop_rel(r_id)
    del presentation.slides._sldIdLst[op.slide_index]


def _apply_move(op: MoveSlideOp, presentation: Presentation) -> None:
    total = len(presentation.slides)
    if op.from_index >= total:
        raise IndexError(f"move_slide from_index out of range: {op.from_index}, total={total}")
    if op.to_index >= total:
        raise IndexError(f"move_slide to_index out of range: {op.to_index}, total={total}")
    if op.from_index == op.to_index:
        return
    slide_id = presentation.slides._sldIdLst[op.from_index]
    presentation.slides._sldIdLst.remove(slide_id)
    presentation.slides._sldIdLst.insert(op.to_index, slide_id)


def _apply_set_slide_size(op: SetSlideSizeOp, presentation: Presentation) -> None:
    if op.preset == "16:9":
        width_inches, height_inches = 13.333, 7.5
    elif op.preset == "4:3":
        width_inches, height_inches = 10.0, 7.5
    else:
        width_inches = float(op.width_inches)  # validated non-null by model
        height_inches = float(op.height_inches)
    presentation.slide_width = Inches(width_inches)
    presentation.slide_height = Inches(height_inches)


def _apply_set_slide_layout(op: SetSlideLayoutOp, presentation: Presentation) -> None:
    if op.slide_index >= len(presentation.slides):
        raise IndexError(
            f"set_slide_layout slide_index out of range: {op.slide_index}, total={len(presentation.slides)}"
        )
    if op.layout_index >= len(presentation.slide_layouts):
        raise IndexError(
            f"set_slide_layout layout_index out of range: {op.layout_index}, total={len(presentation.slide_layouts)}"
        )
    slide = presentation.slides[op.slide_index]
    target_layout = presentation.slide_layouts[op.layout_index]
    for rel in list(slide.part.rels.values()):
        if rel.reltype == _SLIDE_LAYOUT_RELTYPE:
            slide.part.drop_rel(rel.rId)
    slide.part.relate_to(target_layout.part, _SLIDE_LAYOUT_RELTYPE)


def _apply_set_notes(op: SetNotesOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "set_notes")
    slide.notes_slide.notes_text_frame.text = op.text


def _apply_add_textbox(op: AddTextBoxOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "add_textbox")
    shape = slide.shapes.add_textbox(
        Inches(op.x_inches),
        Inches(op.y_inches),
        Inches(op.width_inches),
        Inches(op.height_inches),
    )
    if op.name:
        shape.name = op.name
    _write_text_frame(shape.text_frame, op.text, op.paragraphs, op.vertical_anchor, op.word_wrap)


def _apply_set_shape_text(op: SetShapeTextOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "set_shape_text")
    if op.shape_name is not None:
        shape = _shape_by_name(slide, op.shape_name)
    else:
        shape = _shape_by_index(slide, int(op.shape_index))
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"target shape has no text frame: {shape.name}")
    _write_text_frame(shape.text_frame, op.text, op.paragraphs, op.vertical_anchor, op.word_wrap)


def _apply_add_image(op: AddImageOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "add_image")
    image_path = Path(op.image_path).expanduser().resolve()
    if not image_path.exists():
        raise FileNotFoundError(f"image_path not found: {image_path}")

    x = Inches(op.x_inches)
    y = Inches(op.y_inches)
    box_w = Inches(op.width_inches)
    box_h = Inches(op.height_inches)

    pic = slide.shapes.add_picture(str(image_path), x, y, box_w, box_h)
    iw_px, ih_px = pic.image.size
    image_ratio = iw_px / ih_px if ih_px else 1.0
    box_ratio = box_w / box_h if box_h else 1.0

    if op.fit == "contain":
        if image_ratio > box_ratio:
            new_w = box_w
            new_h = int(box_w / image_ratio)
            pic.left = x
            pic.top = y + int((box_h - new_h) / 2)
        else:
            new_h = box_h
            new_w = int(box_h * image_ratio)
            pic.left = x + int((box_w - new_w) / 2)
            pic.top = y
        pic.width = int(new_w)
        pic.height = int(new_h)
    elif op.fit == "cover":
        if image_ratio > box_ratio:
            target_w = ih_px * box_ratio
            crop = (iw_px - target_w) / (2 * iw_px)
            pic.crop_left = crop
            pic.crop_right = crop
        elif image_ratio < box_ratio:
            target_h = iw_px / box_ratio
            crop = (ih_px - target_h) / (2 * ih_px)
            pic.crop_top = crop
            pic.crop_bottom = crop

    if op.name:
        pic.name = op.name


def _apply_add_shape(op: AddShapeOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "add_shape")
    x = Inches(op.x_inches)
    y = Inches(op.y_inches)
    width = Inches(op.width_inches)
    height = Inches(op.height_inches)

    if op.shape_type == "line":
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x,
            y,
            x + width,
            y + height,
        )
    else:
        shape = slide.shapes.add_shape(_SHAPE_MAP[op.shape_type], x, y, width, height)

    if op.name:
        shape.name = op.name

    if op.fill_color_hex is not None and hasattr(shape, "fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(op.fill_color_hex)

    if op.line_color_hex is not None and hasattr(shape, "line"):
        shape.line.color.rgb = _hex_to_rgb(op.line_color_hex)
    if op.line_width_pt is not None and hasattr(shape, "line"):
        shape.line.width = Pt(op.line_width_pt)

    if op.text and getattr(shape, "has_text_frame", False):
        shape.text_frame.clear()
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = op.text
        if op.font_size_pt is not None:
            paragraph.font.size = Pt(op.font_size_pt)
        if op.text_color_hex is not None:
            paragraph.font.color.rgb = _hex_to_rgb(op.text_color_hex)


def _apply_add_table(op: AddTableOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "add_table")
    rows = len(op.data)
    cols = max(len(row) for row in op.data)
    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(op.x_inches),
        Inches(op.y_inches),
        Inches(op.width_inches),
        Inches(op.height_inches),
    )
    if op.name:
        table_shape.name = op.name
    table = table_shape.table
    for row_idx in range(rows):
        row = op.data[row_idx]
        for col_idx in range(cols):
            text = row[col_idx] if col_idx < len(row) else ""
            cell = table.cell(row_idx, col_idx)
            cell.text_frame.text = text
            paragraph = cell.text_frame.paragraphs[0]
            if op.font_size_pt is not None:
                paragraph.font.size = Pt(op.font_size_pt)
            if op.header and row_idx == 0:
                paragraph.font.bold = True


def _apply_set_slide_background(op: SetSlideBackgroundOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "set_slide_background")
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(op.color_hex)


def _apply_align_shapes(op: AlignShapesOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "align_shapes")
    shapes = [_shape_by_name(slide, name) for name in op.shape_names]
    anchor = shapes[0]

    if op.align in ("left", "center", "right"):
        if op.reference == "slide":
            if op.align == "left":
                target = 0
            elif op.align == "center":
                target = presentation.slide_width / 2
            else:
                target = presentation.slide_width
        else:
            if op.align == "left":
                target = anchor.left
            elif op.align == "center":
                target = anchor.left + anchor.width / 2
            else:
                target = anchor.left + anchor.width
        for shape in shapes:
            if op.align == "left":
                shape.left = int(target)
            elif op.align == "center":
                shape.left = int(target - shape.width / 2)
            else:
                shape.left = int(target - shape.width)
        return

    if op.reference == "slide":
        if op.align == "top":
            target = 0
        elif op.align == "middle":
            target = presentation.slide_height / 2
        else:
            target = presentation.slide_height
    else:
        if op.align == "top":
            target = anchor.top
        elif op.align == "middle":
            target = anchor.top + anchor.height / 2
        else:
            target = anchor.top + anchor.height
    for shape in shapes:
        if op.align == "top":
            shape.top = int(target)
        elif op.align == "middle":
            shape.top = int(target - shape.height / 2)
        else:
            shape.top = int(target - shape.height)


def _apply_distribute_shapes(op: DistributeShapesOp, presentation: Presentation) -> None:
    slide = _slide_or_raise(presentation, op.slide_index, "distribute_shapes")
    shapes = [_shape_by_name(slide, name) for name in op.shape_names]

    if op.direction == "horizontal":
        ordered = sorted(shapes, key=lambda s: s.left)
        start = ordered[0].left
        end = ordered[-1].left + ordered[-1].width
        total_shape_width = sum(shape.width for shape in ordered)
        gap = (end - start - total_shape_width) / (len(ordered) - 1)
        cursor = float(start)
        for shape in ordered:
            shape.left = int(cursor)
            cursor += shape.width + gap
        return

    ordered = sorted(shapes, key=lambda s: s.top)
    start = ordered[0].top
    end = ordered[-1].top + ordered[-1].height
    total_shape_height = sum(shape.height for shape in ordered)
    gap = (end - start - total_shape_height) / (len(ordered) - 1)
    cursor = float(start)
    for shape in ordered:
        shape.top = int(cursor)
        cursor += shape.height + gap


def _to_operations(raw_ops: Iterable[Operation] | list[dict] | dict) -> tuple[list[Operation], OperationPlan | None]:
    if isinstance(raw_ops, dict):
        plan = parse_plan(raw_ops)
        return plan.operations, plan
    raw_list = list(raw_ops)  # type: ignore[arg-type]
    if not raw_list:
        return [], None
    if isinstance(raw_list[0], dict):
        return parse_ops(raw_list), None  # type: ignore[arg-type]
    return raw_list, None  # type: ignore[return-value]


def apply_ops(
    input_pptx: str | Path | None,
    ops: Iterable[Operation] | list[dict] | dict,
    output_pptx: str | Path | None,
    verify: bool = False,
    strict_verify: bool = True,
    template_pptx: str | Path | None = None,
) -> ApplyResult:
    operations, plan = _to_operations(ops)
    template_path_raw = template_pptx if template_pptx is not None else input_pptx
    if template_path_raw is None and plan and plan.template_pptx:
        template_path_raw = plan.template_pptx
    if template_path_raw is None:
        raise ValueError("template_pptx is required")
    if output_pptx is None:
        raise ValueError("output_pptx is required")

    input_path = Path(template_path_raw).expanduser().resolve()
    output_path = Path(output_pptx).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    needs_copy_ops = any(isinstance(op, CopySlideOp) for op in operations)
    copier = None
    if needs_copy_ops:
        SlideCopier, SlideSpec = _import_copy_ops()
        copier = SlideCopier(target_template=input_path, clear_existing=False)
        presentation = copier.presentation
    else:
        SlideSpec = None  # type: ignore[assignment]
        presentation = Presentation(str(input_path))

    for op in operations:
        if isinstance(op, CopySlideOp):
            if copier is None or SlideSpec is None:
                raise RuntimeError("copy_slide requested without initialized copy engine")
            mode = op.mode.value if isinstance(op.mode, CopyMode) else str(op.mode)
            source_path = op.source_path
            if source_path is None:
                if not plan:
                    raise ValueError("copy_slide with reuse_library_index requires plan context")
                if op.reuse_library_index is None:
                    raise ValueError("copy_slide missing source_path and reuse_library_index")
                if op.reuse_library_index >= len(plan.reuse_slide_libraries):
                    raise IndexError(
                        f"reuse_library_index out of range: {op.reuse_library_index}, "
                        f"total={len(plan.reuse_slide_libraries)}"
                    )
                source_path = plan.reuse_slide_libraries[op.reuse_library_index]
            copier.copy_slide(
                SlideSpec(source_path=source_path, slide_index=op.source_slide_index),
                mode=mode,
            )
            continue
        if isinstance(op, CreateSlideOnLayoutOp):
            _apply_create(op, presentation)
            continue
        if isinstance(op, RewriteTextOp):
            _apply_rewrite(op, presentation)
            continue
        if isinstance(op, DeleteSlideOp):
            _apply_delete(op, presentation)
            continue
        if isinstance(op, MoveSlideOp):
            _apply_move(op, presentation)
            continue
        if isinstance(op, SetSlideSizeOp):
            _apply_set_slide_size(op, presentation)
            continue
        if isinstance(op, SetSlideLayoutOp):
            _apply_set_slide_layout(op, presentation)
            continue
        if isinstance(op, SetNotesOp):
            _apply_set_notes(op, presentation)
            continue
        if isinstance(op, AddTextBoxOp):
            _apply_add_textbox(op, presentation)
            continue
        if isinstance(op, SetShapeTextOp):
            _apply_set_shape_text(op, presentation)
            continue
        if isinstance(op, AddImageOp):
            _apply_add_image(op, presentation)
            continue
        if isinstance(op, AddShapeOp):
            _apply_add_shape(op, presentation)
            continue
        if isinstance(op, AddTableOp):
            _apply_add_table(op, presentation)
            continue
        if isinstance(op, SetSlideBackgroundOp):
            _apply_set_slide_background(op, presentation)
            continue
        if isinstance(op, AlignShapesOp):
            _apply_align_shapes(op, presentation)
            continue
        if isinstance(op, DistributeShapesOp):
            _apply_distribute_shapes(op, presentation)
            continue
        raise ValueError(f"Unsupported operation type: {type(op)!r}")

    if copier is not None:
        saved_path = copier.save(output_path)
    else:
        presentation.save(str(output_path))
        saved_path = output_path

    issues: list[str] = []
    if verify:
        report = verify_pptx(saved_path)
        issues = report.issues
        if issues and strict_verify:
            raise ValueError("verification failed: " + "; ".join(issues))

    return ApplyResult(
        output_path=saved_path.resolve(),
        operations_applied=len(operations),
        verify_issues=issues,
    )


def generate_pptx(
    template_pptx: str | Path,
    ops: Iterable[Operation] | list[dict] | dict,
    output_pptx: str | Path,
    verify: bool = False,
    strict_verify: bool = True,
) -> ApplyResult:
    """Generate PPTX from a master/layout template and operation list."""
    return apply_ops(
        input_pptx=None,
        template_pptx=template_pptx,
        ops=ops,
        output_pptx=output_pptx,
        verify=verify,
        strict_verify=strict_verify,
    )
