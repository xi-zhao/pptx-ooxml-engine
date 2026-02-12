from __future__ import annotations

import base64
import importlib.util
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


def _build_source_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "From Source"
    prs.save(str(path))


def _build_target_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Original Title"
    prs.save(str(path))


def _build_three_slide_pptx(path: Path) -> None:
    prs = Presentation()
    titles = ["Slide A", "Slide B", "Slide C"]
    for title in titles:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
    prs.save(str(path))


def _slide_texts(slide) -> list[str]:
    out: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                out.append(text)
    return out


def _write_tiny_png(path: Path) -> None:
    # 1x1 PNG
    raw = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/a5sAAAAASUVORK5CYII="
    )
    path.write_bytes(raw)


def _copy_ops_available() -> bool:
    if importlib.util.find_spec("pptx_copy_ops") is not None:
        return True
    workspace_root = Path(__file__).resolve().parents[2]
    fallback_src = workspace_root / "pptx-copy-ops" / "src"
    return fallback_src.exists()


@pytest.mark.skipif(not _copy_ops_available(), reason="pptx-copy-ops not available")
def test_apply_ops_builds_expected_output(tmp_path: Path) -> None:
    from pptx_ooxml_engine.engine import apply_ops

    source = tmp_path / "source.pptx"
    target = tmp_path / "target.pptx"
    output = tmp_path / "output.pptx"
    _build_source_pptx(source)
    _build_target_pptx(target)

    result = apply_ops(
        input_pptx=target,
        ops=[
            {"op": "rewrite_text", "slide_index": 0, "find": "Original", "replace": "Rewritten"},
            {"op": "copy_slide", "source_path": str(source), "source_slide_index": 0, "mode": "shape"},
            {"op": "create_slide_on_layout", "layout_index": 0, "title": "Created Slide", "body": "Body"},
        ],
        output_pptx=output,
        verify=True,
    )

    assert result.output_path == output.resolve()
    prs = Presentation(str(output))
    assert len(prs.slides) == 3

    s1_texts = _slide_texts(prs.slides[0])
    s2_texts = _slide_texts(prs.slides[1])
    s3_texts = _slide_texts(prs.slides[2])
    assert any("Rewritten Title" == t for t in s1_texts)
    assert any("From Source" == t for t in s2_texts)
    assert any("Created Slide" == t for t in s3_texts)


def test_apply_ops_can_take_template_from_plan_dict(tmp_path: Path) -> None:
    from pptx_ooxml_engine.engine import apply_ops

    template = tmp_path / "template.pptx"
    output = tmp_path / "output_from_plan.pptx"
    _build_target_pptx(template)

    result = apply_ops(
        input_pptx=None,
        ops={
            "template_pptx": str(template),
            "reuse_slide_libraries": ["/tmp/reuse1.pptx"],
            "operations": [
                {"op": "rewrite_text", "slide_index": 0, "find": "Original", "replace": "Plan"},
                {"op": "create_slide_on_layout", "layout_index": 0, "title": "Added By Plan"},
            ],
        },
        output_pptx=output,
        verify=True,
    )

    assert result.output_path == output.resolve()
    prs = Presentation(str(output))
    assert len(prs.slides) == 2
    assert any("Plan Title" == t for t in _slide_texts(prs.slides[0]))


@pytest.mark.skipif(not _copy_ops_available(), reason="pptx-copy-ops not available")
def test_apply_ops_copy_from_reuse_library_index(tmp_path: Path) -> None:
    from pptx_ooxml_engine.engine import apply_ops

    template = tmp_path / "template.pptx"
    reuse_lib = tmp_path / "reuse_lib.pptx"
    output = tmp_path / "output_from_reuse_index.pptx"
    _build_target_pptx(template)
    _build_source_pptx(reuse_lib)

    result = apply_ops(
        input_pptx=None,
        ops={
            "template_pptx": str(template),
            "reuse_slide_libraries": [str(reuse_lib)],
            "operations": [
                {"op": "copy_slide", "reuse_library_index": 0, "source_slide_index": 0, "mode": "shape"},
            ],
        },
        output_pptx=output,
        verify=True,
    )

    assert result.output_path == output.resolve()
    prs = Presentation(str(output))
    assert len(prs.slides) == 2
    assert any("From Source" == t for t in _slide_texts(prs.slides[1]))


def test_apply_ops_delete_slide_then_create(tmp_path: Path) -> None:
    from pptx_ooxml_engine.engine import apply_ops

    template = tmp_path / "template.pptx"
    output = tmp_path / "output_delete_then_create.pptx"
    _build_target_pptx(template)  # contains one initial slide

    result = apply_ops(
        input_pptx=None,
        ops={
            "template_pptx": str(template),
            "operations": [
                {"op": "delete_slide", "slide_index": 0},
                {"op": "create_slide_on_layout", "layout_index": 0, "title": "Only Slide"},
            ],
        },
        output_pptx=output,
        verify=True,
    )

    assert result.output_path == output.resolve()
    prs = Presentation(str(output))
    assert len(prs.slides) == 1
    assert any("Only Slide" == t for t in _slide_texts(prs.slides[0]))


def test_apply_ops_structure_management_ops(tmp_path: Path) -> None:
    from pptx_ooxml_engine.engine import apply_ops

    template = tmp_path / "template_three_slides.pptx"
    output = tmp_path / "output_structure_ops.pptx"
    _build_three_slide_pptx(template)

    result = apply_ops(
        input_pptx=None,
        ops={
            "template_pptx": str(template),
            "operations": [
                {"op": "move_slide", "from_index": 0, "to_index": 2},
                {"op": "set_slide_layout", "slide_index": 0, "layout_index": 1},
                {"op": "set_notes", "slide_index": 0, "text": "Speaker Notes"},
                {"op": "set_slide_size", "preset": "4:3"},
            ],
        },
        output_pptx=output,
        verify=True,
    )

    assert result.output_path == output.resolve()
    prs = Presentation(str(output))
    assert len(prs.slides) == 3
    # move_slide: A B C -> B C A
    assert any("Slide B" == t for t in _slide_texts(prs.slides[0]))
    assert any("Slide C" == t for t in _slide_texts(prs.slides[1]))
    assert any("Slide A" == t for t in _slide_texts(prs.slides[2]))

    # set_slide_size preset 4:3
    assert prs.slide_width == Inches(10)
    assert prs.slide_height == Inches(7.5)

    # set_slide_layout: slide0 layout should be layout 1
    assert str(prs.slides[0].slide_layout.part.partname) == str(prs.slide_layouts[1].part.partname)

    # set_notes
    assert "Speaker Notes" in prs.slides[0].notes_slide.notes_text_frame.text


def test_apply_ops_without_copy_does_not_import_copy_module(tmp_path: Path, monkeypatch) -> None:
    import pptx_ooxml_engine.engine as engine

    template = tmp_path / "template_no_copy.pptx"
    output = tmp_path / "output_no_copy.pptx"
    _build_target_pptx(template)

    def _boom():  # pragma: no cover - assertion path
        raise AssertionError("copy ops import should be lazy")

    monkeypatch.setattr(engine, "_import_copy_ops", _boom)

    result = engine.apply_ops(
        input_pptx=None,
        ops={
            "template_pptx": str(template),
            "operations": [
                {"op": "rewrite_text", "slide_index": 0, "find": "Original", "replace": "Lazy"},
                {"op": "set_notes", "slide_index": 0, "text": "No copy used"},
            ],
        },
        output_pptx=output,
        verify=True,
    )

    assert result.output_path == output.resolve()
    prs = Presentation(str(output))
    assert any("Lazy Title" == t for t in _slide_texts(prs.slides[0]))


def test_apply_ops_v1_content_and_layout_ops(tmp_path: Path) -> None:
    from pptx_ooxml_engine.engine import apply_ops

    template = tmp_path / "template_v1_ops.pptx"
    output = tmp_path / "output_v1_ops.pptx"
    image_path = tmp_path / "tiny.png"
    _build_target_pptx(template)
    _write_tiny_png(image_path)

    result = apply_ops(
        input_pptx=None,
        ops={
            "template_pptx": str(template),
            "operations": [
                {
                    "op": "add_shape",
                    "slide_index": 0,
                    "shape_type": "rect",
                    "x_inches": 1.0,
                    "y_inches": 1.2,
                    "width_inches": 2.0,
                    "height_inches": 1.0,
                    "text": "Box 1",
                    "name": "box1",
                    "fill_color_hex": "0A84FF",
                },
                {
                    "op": "add_shape",
                    "slide_index": 0,
                    "shape_type": "rect",
                    "x_inches": 4.0,
                    "y_inches": 2.0,
                    "width_inches": 2.0,
                    "height_inches": 1.0,
                    "text": "Box 2",
                    "name": "box2",
                },
                {
                    "op": "add_shape",
                    "slide_index": 0,
                    "shape_type": "rect",
                    "x_inches": 7.2,
                    "y_inches": 2.8,
                    "width_inches": 2.0,
                    "height_inches": 1.0,
                    "text": "Box 3",
                    "name": "box3",
                },
                {
                    "op": "align_shapes",
                    "slide_index": 0,
                    "shape_names": ["box1", "box2", "box3"],
                    "align": "top",
                    "reference": "first",
                },
                {
                    "op": "distribute_shapes",
                    "slide_index": 0,
                    "shape_names": ["box1", "box2", "box3"],
                    "direction": "horizontal",
                },
                {
                    "op": "set_shape_text",
                    "slide_index": 0,
                    "shape_name": "box2",
                    "paragraphs": [
                        {
                            "text": "Updated Box 2",
                            "font_size_pt": 18,
                            "bold": True,
                            "alignment": "center",
                        },
                        {
                            "text": "bullet line",
                            "list_type": "bullet",
                            "level": 1,
                        },
                        {
                            "text": "number line",
                            "list_type": "number",
                            "level": 1,
                        },
                    ],
                },
                {
                    "op": "add_textbox",
                    "slide_index": 0,
                    "x_inches": 0.8,
                    "y_inches": 4.2,
                    "width_inches": 3.5,
                    "height_inches": 2.0,
                    "name": "briefing",
                    "paragraphs": [
                        {"text": "Roadmap", "font_size_pt": 20, "bold": True},
                        {"text": "item one", "list_type": "bullet", "level": 1},
                        {"text": "item two", "list_type": "number", "level": 1},
                    ],
                },
                {
                    "op": "add_table",
                    "slide_index": 0,
                    "x_inches": 4.5,
                    "y_inches": 4.2,
                    "width_inches": 4.5,
                    "height_inches": 2.0,
                    "data": [["Metric", "Value"], ["Revenue", "1148"]],
                    "header": True,
                },
                {
                    "op": "add_image",
                    "slide_index": 0,
                    "image_path": str(image_path),
                    "x_inches": 10.2,
                    "y_inches": 0.3,
                    "width_inches": 2.2,
                    "height_inches": 1.1,
                    "fit": "cover",
                },
                {
                    "op": "set_slide_background",
                    "slide_index": 0,
                    "color_hex": "112233",
                },
            ],
        },
        output_pptx=output,
        verify=True,
    )

    assert result.output_path == output.resolve()
    prs = Presentation(str(output))
    slide = prs.slides[0]

    box1 = next(shape for shape in slide.shapes if shape.name == "box1")
    box2 = next(shape for shape in slide.shapes if shape.name == "box2")
    box3 = next(shape for shape in slide.shapes if shape.name == "box3")
    assert box1.top == box2.top == box3.top
    gap1 = box2.left - (box1.left + box1.width)
    gap2 = box3.left - (box2.left + box2.width)
    assert abs(gap1 - gap2) <= 2

    assert "Updated Box 2" in box2.text
    briefing = next(shape for shape in slide.shapes if shape.name == "briefing")
    assert "Roadmap" in briefing.text
    assert "item one" in briefing.text

    picture_count = sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert picture_count == 1

    table_shapes = [shape for shape in slide.shapes if shape.has_table]
    assert len(table_shapes) == 1
    assert table_shapes[0].table.cell(0, 0).text_frame.text == "Metric"
    assert table_shapes[0].table.cell(1, 1).text_frame.text == "1148"

    fill = slide.background.fill
    assert fill.fore_color.rgb == RGBColor.from_string("112233")
