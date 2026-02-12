from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def test_generate_pptx_alias_matches_apply_ops(tmp_path: Path) -> None:
    from pptx_ooxml_engine import apply_ops, generate_pptx

    source = tmp_path / "template.pptx"
    out_a = tmp_path / "a.pptx"
    out_b = tmp_path / "b.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Old"
    prs.save(str(source))

    ops = [{"op": "rewrite_text", "slide_index": 0, "find": "Old", "replace": "New"}]
    result_a = apply_ops(source, ops, out_a, verify=True)
    result_b = generate_pptx(template_pptx=source, ops=ops, output_pptx=out_b, verify=True)

    assert result_a.operations_applied == result_b.operations_applied == 1
    assert result_a.verify_issues == result_b.verify_issues == []
