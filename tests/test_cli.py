from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation


def test_cli_runs_ops_file(tmp_path: Path) -> None:
    project_root = Path(__file__).resolve().parents[1]
    input_pptx = tmp_path / "template.pptx"
    output_pptx = tmp_path / "output.pptx"
    ops_file = tmp_path / "ops.json"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Old Title"
    prs.save(str(input_pptx))

    ops_file.write_text(
        json.dumps(
            {
                "operations": [
                    {
                        "op": "rewrite_text",
                        "slide_index": 0,
                        "find": "Old",
                        "replace": "New",
                    },
                    {
                        "op": "create_slide_on_layout",
                        "layout_index": 0,
                        "title": "Added",
                    },
                ]
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    completed = subprocess.run(
        [
                sys.executable,
                "-m",
                "pptx_ooxml_engine.cli",
                "--template",
                str(input_pptx),
                "--ops-file",
                str(ops_file),
            "--output",
            str(output_pptx),
            "--verify",
        ],
        cwd=str(project_root),
        env={"PYTHONPATH": str(project_root / "src")},
        capture_output=True,
        text=True,
        check=False,
    )

    assert completed.returncode == 0, completed.stderr
    assert output_pptx.exists()


def test_cli_can_use_template_from_ops_file(tmp_path: Path) -> None:
    project_root = Path(__file__).resolve().parents[1]
    template_pptx = tmp_path / "template.pptx"
    output_pptx = tmp_path / "output_from_ops_template.pptx"
    ops_file = tmp_path / "ops_with_template.json"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Old Title"
    prs.save(str(template_pptx))

    ops_file.write_text(
        json.dumps(
            {
                "template_pptx": str(template_pptx),
                "operations": [
                    {
                        "op": "rewrite_text",
                        "slide_index": 0,
                        "find": "Old",
                        "replace": "New",
                    }
                ],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    completed = subprocess.run(
        [
            sys.executable,
            "-m",
            "pptx_ooxml_engine.cli",
            "--ops-file",
            str(ops_file),
            "--output",
            str(output_pptx),
            "--verify",
        ],
        cwd=str(project_root),
        env={"PYTHONPATH": str(project_root / "src")},
        capture_output=True,
        text=True,
        check=False,
    )

    assert completed.returncode == 0, completed.stderr
    assert output_pptx.exists()
